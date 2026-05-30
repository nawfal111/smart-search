#!/usr/bin/env python3
"""
Smart Search — PowerPoint Presentation Generator
Run: python3 generate_pptx.py
Output: smart_search_presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

DIR  = os.path.dirname(os.path.abspath(__file__))
W    = Inches(13.33)
H    = Inches(7.5)

# ── Palette ────────────────────────────────────────────────────────────────────
BG     = RGBColor(0x0D, 0x13, 0x23)   # near-black navy  – background
PANEL  = RGBColor(0x19, 0x24, 0x3A)   # dark slate       – cards
BORDER = RGBColor(0x2D, 0x3F, 0x58)   # muted blue       – borders
BLUE   = RGBColor(0x38, 0x7E, 0xF5)   # blue             – extension / primary
GREEN  = RGBColor(0x0D, 0xB5, 0x7C)   # emerald          – local / success
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)   # amber            – AI / cloud
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)   # violet           – Pinecone / vectors
RED    = RGBColor(0xF4, 0x3F, 0x5E)   # rose             – problem / warning
WHITE  = RGBColor(0xF8, 0xFA, 0xFC)   # near-white       – primary text
GREY   = RGBColor(0x94, 0xA3, 0xB8)   # slate grey       – secondary text
DIM    = RGBColor(0x3D, 0x50, 0x6A)   # dim blue         – subtle elements


# ── Core helpers ───────────────────────────────────────────────────────────────

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide


def box(slide, x, y, w, h, fill, border=None, bw=1.0, rounding=False):
    """Solid rectangle, optionally with border and rounding."""
    t = 5 if rounding else 1   # 5=RoundedRect, 1=Rectangle
    s = slide.shapes.add_shape(t, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    if rounding:
        # keep rounding small
        try:
            s.adjustments[0] = 0.08
        except Exception:
            pass
    return s


def circle(slide, cx, cy, r, fill, border=None):
    """Circle centred at (cx, cy) with radius r."""
    s = slide.shapes.add_shape(9, cx - r, cy - r, 2 * r, 2 * r)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def line_shape(slide, x1, y1, x2, y2, color, width=1.5):
    """Draw a straight line connector."""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def t(slide, text, x, y, w, h,
      size=18, bold=False, italic=False,
      color=None, align=PP_ALIGN.LEFT, wrap=True):
    """Single-run text box."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.italic = italic
    run.font.size   = Pt(size)
    run.font.color.rgb = color or WHITE
    return tb


def tlines(slide, entries, x, y, w, h, wrap=True, spacing=1.1):
    """
    Multi-paragraph text box.
    entries = [(text, size, bold, color, align), ...]
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, (text, size, bold, color, align) in enumerate(entries):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = text
        run.font.bold  = bold
        run.font.size  = Pt(size)
        run.font.color.rgb = color or WHITE
    return tb


def morph(slide):
    """Attach a Morph transition to the slide."""
    xml = (
        '<p:transition'
        ' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"'
        ' spd="med" advTm="0">'
        '<p14:transition inverted="0">'
        '<p14:morph option="byObject"/>'
        '</p14:transition>'
        '</p:transition>'
    )
    slide._element.append(etree.fromstring(xml))


def top_bar(slide, color, h=Inches(0.055)):
    box(slide, 0, 0, W, h, color)


def slide_label(slide, text, color):
    """Small coloured pill in top-left."""
    b = box(slide, Inches(0.42), Inches(0.15), Inches(2.0), Inches(0.38),
            color, rounding=True)
    t(slide, text, Inches(0.42), Inches(0.15), Inches(2.0), Inches(0.38),
      size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def section_title(slide, text, color=WHITE, size=36, y=Inches(0.6)):
    t(slide, text, Inches(0.55), y, Inches(12.0), Inches(0.9),
      size=size, bold=True, color=color)


def divider(slide, y, color=DIM):
    box(slide, Inches(0.55), y, Inches(12.2), Inches(0.018), color)


# ── Card helper ────────────────────────────────────────────────────────────────

def card(slide, x, y, w, h, accent_color, title, body_lines,
         title_size=13, body_size=10.5, padding=Inches(0.18)):
    """Rounded dark card with coloured left border."""
    box(slide, x, y, w, h, PANEL, BORDER, bw=0.8, rounding=True)
    box(slide, x, y, Inches(0.055), h, accent_color)            # left accent stripe
    t(slide, title,
      x + Inches(0.2), y + padding * 0.6, w - Inches(0.3), Inches(0.4),
      size=title_size, bold=True, color=WHITE)
    body = "\n".join(body_lines)
    t(slide, body,
      x + Inches(0.2), y + padding * 0.6 + Inches(0.38),
      w - Inches(0.3), h - Inches(0.6),
      size=body_size, color=GREY, wrap=True)


def step_bubble(slide, cx, cy, r, number, color):
    """Numbered circle for pipeline steps."""
    circle(slide, cx, cy, r, color)
    t(slide, str(number),
      cx - r, cy - r, 2 * r, 2 * r,
      size=int(r.inches * 55), bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def arrow_right(slide, x, y, length, color=DIM, width=2.0):
    line_shape(slide, x, y, x + length, y, color, width)
    # Arrow head (small triangle using a thin line)
    aw = Inches(0.12)
    line_shape(slide, x + length, y, x + length - aw, y - aw * 0.6, color, width)
    line_shape(slide, x + length, y, x + length - aw, y + aw * 0.6, color, width)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    s = new_slide(prs)
    morph(s)

    # Large geometric accent — blue rectangle top-right
    box(s, Inches(9.5), Inches(-0.5), Inches(5.0), Inches(5.0),
        RGBColor(0x1A, 0x35, 0x6E), rounding=False)
    box(s, Inches(10.8), Inches(0.4), Inches(3.5), Inches(3.5),
        RGBColor(0x22, 0x44, 0x88), rounding=False)

    # Accent dots — bottom-left
    for i, c in enumerate([BLUE, GREEN, AMBER]):
        circle(s, Inches(0.55 + i * 0.52), Inches(7.15), Inches(0.13), c)

    # Tag line above title
    t(s, "MASTER'S RESEARCH PROJECT",
      Inches(0.55), Inches(1.6), Inches(8.0), Inches(0.5),
      size=11, bold=True, color=BLUE, align=PP_ALIGN.LEFT)

    # Main title
    t(s, "Smart Search",
      Inches(0.55), Inches(2.05), Inches(9.0), Inches(1.8),
      size=72, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    t(s, "Semantic Code Search Engine for VS Code",
      Inches(0.55), Inches(3.75), Inches(9.5), Inches(0.7),
      size=22, bold=False, color=GREY, align=PP_ALIGN.LEFT)

    # Tagline
    t(s, "Find code by what it does — not what it's called",
      Inches(0.55), Inches(4.4), Inches(9.5), Inches(0.6),
      size=16, italic=True, color=DIM, align=PP_ALIGN.LEFT)

    # Bottom divider
    box(s, Inches(0.55), Inches(5.5), Inches(1.8), Inches(0.04), BLUE)

    # Author + date
    tlines(s, [
        ("Nawfal Jalloul", 13, True,  WHITE,  PP_ALIGN.LEFT),
        ("May 2026",        11, False, GREY,   PP_ALIGN.LEFT),
    ], Inches(0.55), Inches(5.65), Inches(5.0), Inches(0.9))


def slide_02_problem(prs):
    s = new_slide(prs)
    morph(s)
    top_bar(s, RED)
    slide_label(s, "THE PROBLEM", RED)

    section_title(s, "Traditional Search is Broken", color=WHITE, size=30, y=Inches(0.65))
    divider(s, Inches(1.5))

    # Left panel — what you type
    box(s, Inches(0.45), Inches(1.65), Inches(5.8), Inches(4.9), PANEL, BORDER, rounding=True)
    t(s, "You search for...",
      Inches(0.7), Inches(1.8), Inches(5.2), Inches(0.45),
      size=11, color=GREY, bold=False)
    t(s, '"fetch product\nfrom database"',
      Inches(0.7), Inches(2.22), Inches(5.2), Inches(1.2),
      size=20, bold=True, color=AMBER, align=PP_ALIGN.LEFT)

    # Big 0 results
    t(s, "0",
      Inches(1.5), Inches(3.55), Inches(2.2), Inches(1.8),
      size=110, bold=True, color=RED, align=PP_ALIGN.CENTER)
    t(s, "RESULTS",
      Inches(1.5), Inches(5.05), Inches(2.2), Inches(0.5),
      size=14, bold=True, color=RED, align=PP_ALIGN.CENTER)

    # Right panel — why
    box(s, Inches(6.65), Inches(1.65), Inches(6.2), Inches(4.9), PANEL, BORDER, rounding=True)
    t(s, "The function exists — it's just named differently:",
      Inches(6.9), Inches(1.82), Inches(5.8), Inches(0.5),
      size=11, color=GREY)

    fn_names = [
        ("getProductInfo()",    BLUE),
        ("retrieveItem()",      GREEN),
        ("fetchFromDatabase()", AMBER),
        ("loadProductById()",   PURPLE),
    ]
    for i, (name, color) in enumerate(fn_names):
        yy = Inches(2.42 + i * 0.82)
        box(s, Inches(6.9), yy, Inches(5.7), Inches(0.6), BORDER, rounding=True)
        t(s, name, Inches(7.05), yy + Inches(0.04), Inches(5.3), Inches(0.55),
          size=15, bold=True, color=color, align=PP_ALIGN.LEFT)

    t(s, "Traditional search: 0 matches.   Smart Search: finds them all.",
      Inches(0.45), Inches(6.72), Inches(12.3), Inches(0.45),
      size=11, italic=True, color=DIM, align=PP_ALIGN.CENTER)


def slide_03_solution(prs):
    s = new_slide(prs)
    morph(s)
    top_bar(s, BLUE)
    slide_label(s, "THE SOLUTION", BLUE)

    section_title(s, "Search by Meaning, Not Words", size=30, y=Inches(0.65))
    divider(s, Inches(1.5))

    # Central flow: Query → Embed → Match
    centres = [Inches(2.2), Inches(6.65), Inches(11.1)]
    labels  = ["Natural Language\nQuery", "Semantic\nVector Space", "Code\nFunction"]
    colors  = [BLUE, PURPLE, GREEN]
    icons   = ['"fetch product\nfrom database"', "1,536-float\nvector", "getProductInfo()"]

    for i, (cx, label, color, icon) in enumerate(zip(centres, labels, colors, icons)):
        # Circle
        r = Inches(1.35)
        circle(s, cx, Inches(3.6), r, color)
        t(s, icon,
          cx - r, Inches(3.6) - r, 2*r, 2*r,
          size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Label below
        t(s, label, cx - Inches(1.5), Inches(5.15), Inches(3.0), Inches(0.7),
          size=12, color=GREY, align=PP_ALIGN.CENTER)

        # Arrow between circles
        if i < 2:
            ax = cx + r + Inches(0.08)
            next_cx = centres[i+1]
            arrow_right(s, ax, Inches(3.6), next_cx - r - Inches(0.08) - ax,
                       BORDER, 2.5)
            mid_x = (ax + next_cx - r) / 2
            step_labels = ["embed via\nVoyage AI", "cosine\nsimilarity\n(HNSW)"]
            t(s, step_labels[i], mid_x - Inches(0.9), Inches(3.0), Inches(1.8), Inches(0.65),
              size=9.5, color=AMBER, align=PP_ALIGN.CENTER, italic=True)

    # Key insight box at bottom
    box(s, Inches(2.5), Inches(6.1), Inches(8.3), Inches(0.85), PANEL, BLUE, bw=1.0, rounding=True)
    t(s, '82% similarity match — even with zero shared words between query and function name',
      Inches(2.7), Inches(6.22), Inches(7.9), Inches(0.6),
      size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_04_architecture(prs):
    s = new_slide(prs)
    morph(s)
    top_bar(s, PURPLE)
    slide_label(s, "ARCHITECTURE", PURPLE)
    section_title(s, "Three-Layer Design", size=28, y=Inches(0.65))
    divider(s, Inches(1.46))

    cols = [
        (Inches(0.42),  Inches(4.15), BLUE,   "VS Code Extension",     "TypeScript / Node.js",
         ["File watcher", "Code chunker (12 langs)", "Hash store — index.json",
          "Two-phase batch indexer", "Search UI webview", "Status bar"]),
        (Inches(4.72),  Inches(4.15), GREEN,  "Python Backend",         "localhost:8000",
         ["ThreadingHTTPServer", "GPT summarizer (parallel)", "Voyage AI embedder",
          "Pinecone client", "Normal search (regex)", "/config · /health · /done"]),
        (Inches(9.02),  Inches(4.1),  AMBER,  "Cloud Services",         "External APIs",
         ["OpenAI  GPT-4o-mini", "Voyage AI  voyage-code-2", "Pinecone  (HNSW)", "",
          "namespace: {projectId}::{userId}", ""]),
    ]

    for x, cw, color, title, subtitle, items in cols:
        box(s, x, Inches(1.6), cw, Inches(5.5), PANEL, BORDER, bw=0.7, rounding=True)
        box(s, x, Inches(1.6), cw, Inches(0.06), color)  # top accent
        t(s, title,    x + Inches(0.18), Inches(1.72), cw - Inches(0.25), Inches(0.48),
          size=15, bold=True, color=color)
        t(s, subtitle, x + Inches(0.18), Inches(2.15), cw - Inches(0.25), Inches(0.32),
          size=10, color=GREY)
        box(s, x + Inches(0.1), Inches(2.45), cw - Inches(0.2), Inches(0.018), DIM)
        for i, item in enumerate(items):
            if item:
                t(s, f"• {item}", x + Inches(0.22), Inches(2.57 + i * 0.55),
                  cw - Inches(0.32), Inches(0.5),
                  size=10.5, color=WHITE if i == 0 else GREY)

    # HTTP arrows between columns
    for ax, lbl in [(Inches(4.57), "HTTP"), (Inches(8.87), "API")]:
        arrow_right(s, ax, Inches(4.35), Inches(0.15), BORDER, 1.5)
        t(s, lbl, ax - Inches(0.05), Inches(4.05), Inches(0.3), Inches(0.3),
          size=8, color=DIM, align=PP_ALIGN.CENTER)


def slide_05_indexing(prs):
    s = new_slide(prs)
    morph(s)
    top_bar(s, GREEN)
    slide_label(s, "INDEXING PIPELINE", GREEN)
    section_title(s, "From File Save → Pinecone Vector", size=28, y=Inches(0.65))
    divider(s, Inches(1.46))

    steps = [
        ("1", "File\nSaved",        BG,     GREEN,  "onDidSave\nevent"),
        ("2", "Hash\nCheck",        BG,     BLUE,   "MD5 file\n+ function"),
        ("3", "Chunk\nFunctions",   BG,     BLUE,   "brace/indent\n12 langs"),
        ("4", "GPT\nSummary",       BG,     AMBER,  "parallel\n50 at once"),
        ("5", "Embed\n1,536 floats",BG,     AMBER,  "voyage-\ncode-2"),
        ("6", "Pinecone\nUpsert",   BG,     PURPLE, "namespace\nisolation"),
    ]

    step_w = Inches(1.75)
    gap    = Inches(0.32)
    start_x= Inches(0.45)
    cy     = Inches(4.0)
    r      = Inches(0.62)

    for i, (num, label, bg, color, sub) in enumerate(steps):
        cx = start_x + i * (step_w + gap) + step_w / 2

        # Circle
        circle(s, cx, cy, r, color)
        t(s, num, cx - r, cy - r, 2*r, 2*r,
          size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Step label below circle
        t(s, label, cx - step_w/2 - Inches(0.05), cy + r + Inches(0.1),
          step_w + Inches(0.1), Inches(0.7),
          size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Sub-label
        t(s, sub, cx - step_w/2, cy + r + Inches(0.85),
          step_w, Inches(0.55),
          size=9, color=GREY, align=PP_ALIGN.CENTER)

        # Horizontal connecting line
        if i < len(steps) - 1:
            lx = cx + r + Inches(0.06)
            next_cx = start_x + (i+1)*(step_w+gap) + step_w/2
            arrow_right(s, lx, cy, next_cx - r - Inches(0.06) - lx, BORDER, 2.0)

    # Crash-safe note
    box(s, Inches(0.45), Inches(6.35), Inches(12.3), Inches(0.72), PANEL, GREEN, bw=0.8, rounding=True)
    t(s, "🔒  Crash-safe: index.json saved after every batch — if interrupted, next run resumes from last checkpoint",
      Inches(0.65), Inches(6.47), Inches(11.8), Inches(0.5),
      size=11, color=WHITE)

    # Phase labels
    box(s, Inches(0.45), Inches(2.1), Inches(0.7), Inches(1.8), PANEL, rounding=True)
    t(s, "SCAN", Inches(0.45), Inches(2.1), Inches(0.7), Inches(1.8),
      size=8, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    box(s, Inches(7.0), Inches(2.1), Inches(0.7), Inches(1.8), PANEL, rounding=True)
    t(s, "AI", Inches(7.0), Inches(2.1), Inches(0.7), Inches(1.8),
      size=8, bold=True, color=AMBER, align=PP_ALIGN.CENTER)


def slide_06_summaries(prs):
    s = new_slide(prs)
    morph(s)
    top_bar(s, AMBER)
    slide_label(s, "KEY INNOVATION", AMBER)
    section_title(s, "GPT Summaries Bridge the Vocabulary Gap", size=26, y=Inches(0.65))
    divider(s, Inches(1.46))

    # Code box — left
    box(s, Inches(0.42), Inches(1.65), Inches(4.5), Inches(2.8), PANEL, BORDER, rounding=True)
    t(s, "CODE", Inches(0.65), Inches(1.78), Inches(2.0), Inches(0.35),
      size=9, bold=True, color=BLUE)
    code = ("public function getProductInfo($id) {\n"
            "  $sql = \"SELECT * FROM products\n"
            "          WHERE id = ?\";\n"
            "  $result = $db->query($sql, [$id]);\n"
            "  error_log(\"Fetching: \" . $id);\n"
            "  return $result->fetch_assoc();\n"
            "}")
    t(s, code, Inches(0.55), Inches(2.1), Inches(4.25), Inches(1.75),
      size=9.5, color=GREEN, italic=False)

    # GPT arrow
    t(s, "GPT-4o-mini\nsummarizes",
      Inches(5.1), Inches(2.4), Inches(1.45), Inches(0.8),
      size=10.5, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    arrow_right(s, Inches(4.95), Inches(2.95), Inches(1.6), AMBER, 2.5)

    # Summary box — right of code
    box(s, Inches(6.72), Inches(1.65), Inches(6.0), Inches(2.8), PANEL, AMBER, bw=1.0, rounding=True)
    t(s, "GPT PLAIN-ENGLISH SUMMARY", Inches(6.95), Inches(1.78), Inches(5.5), Inches(0.35),
      size=9, bold=True, color=AMBER)
    summary = ('"Fetches a single product\'s full details from the\n'
               'database by its numeric ID. Executes a parameterized\n'
               'SQL SELECT query on the products table.\n'
               'Includes error logging of the requested product ID."')
    t(s, summary, Inches(6.95), Inches(2.14), Inches(5.5), Inches(1.7),
      size=11, italic=True, color=WHITE)

    # Impact table
    box(s, Inches(0.42), Inches(4.7), Inches(12.5), Inches(2.45), PANEL, BORDER, rounding=True)
    t(s, "EMBEDDING SCORE IMPACT", Inches(0.65), Inches(4.82), Inches(5.0), Inches(0.35),
      size=9, bold=True, color=GREY)

    rows = [
        ('Query',                              'Without Summary', 'With Summary', 'Gain'),
        ('"fetch product from database by id"', '46%',            '82%',          '+36%'),
        ('"check if user is authenticated"',    '38%',            '71%',          '+33%'),
        ('"write product data to database"',    '41%',            '68%',          '+27%'),
        ('"log error to file"',                 '52%',            '74%',          '+22%'),
    ]
    col_x = [Inches(0.6), Inches(6.9), Inches(9.3), Inches(11.4)]
    col_w = [Inches(6.1), Inches(2.2), Inches(2.2), Inches(1.5)]
    for ri, row in enumerate(rows):
        yy = Inches(5.15 + ri * 0.38)
        if ri == 0:
            for ci, cell in enumerate(row):
                t(s, cell, col_x[ci], yy, col_w[ci], Inches(0.36),
                  size=9, bold=True, color=GREY, align=PP_ALIGN.LEFT)
        else:
            colors_row = [WHITE, GREY, GREEN, AMBER]
            for ci, (cell, cc) in enumerate(zip(row, colors_row)):
                t(s, cell, col_x[ci], yy, col_w[ci], Inches(0.36),
                  size=10.5, bold=(ci > 1), color=cc, align=PP_ALIGN.LEFT)
        if ri == 0:
            box(s, Inches(0.55), Inches(5.5), Inches(12.3), Inches(0.02), DIM)


def slide_07_search(prs):
    s = new_slide(prs)
    morph(s)
    top_bar(s, PURPLE)
    slide_label(s, "AI SEARCH PIPELINE", PURPLE)
    section_title(s, "How a Query Becomes a Result", size=28, y=Inches(0.65))
    divider(s, Inches(1.46))

    steps = [
        ("User\nQuery",        BLUE,   '"fetch product\nfrom DB"'),
        ("Voyage AI\nEmbed",   AMBER,  "1,536-float\nvector"),
        ("Pinecone\nHNSW",     PURPLE, "top-10\nnearest"),
        ("Threshold\nFilter",  GREEN,  "≥ 35%\nscore"),
        ("GPT\nLine Locator",  AMBER,  "exact line\nper result"),
        ("Open\nFile",         GREEN,  "line 47\nhighlighted"),
    ]

    step_w = Inches(1.72)
    gap    = Inches(0.28)
    sy     = Inches(3.2)
    bh     = Inches(1.75)

    for i, (label, color, sub) in enumerate(steps):
        x = Inches(0.42) + i * (step_w + gap)
        box(s, x, sy, step_w, bh, PANEL, BORDER, rounding=True)
        box(s, x, sy, step_w, Inches(0.055), color)
        t(s, label, x + Inches(0.1), sy + Inches(0.12), step_w - Inches(0.15), Inches(0.55),
          size=11.5, bold=True, color=color, align=PP_ALIGN.CENTER)
        t(s, sub,   x + Inches(0.1), sy + Inches(0.68), step_w - Inches(0.15), Inches(0.9),
          size=10,  color=GREY, align=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            ax = x + step_w + Inches(0.04)
            arrow_right(s, ax, sy + bh/2, gap - Inches(0.08), BORDER, 2.0)

    # Bottom note boxes
    notes = [
        (BLUE,   "Both layers validate query length\n(frontend + backend safety net)"),
        (PURPLE, "HNSW approximate nearest-neighbour\nsub-millisecond at any scale"),
        (AMBER,  "All GPT line-locator calls fire\nin parallel — total ≈ one GPT call"),
    ]
    nw = Inches(3.8)
    for i, (color, note) in enumerate(notes):
        nx = Inches(0.42) + i * (nw + Inches(0.43))
        box(s, nx, Inches(5.35), nw, Inches(1.6), PANEL, color, bw=0.7, rounding=True)
        t(s, note, nx + Inches(0.18), Inches(5.48), nw - Inches(0.3), Inches(1.3),
          size=10.5, color=WHITE)


def slide_08_vectors(prs):
    s = new_slide(prs)
    morph(s)
    top_bar(s, PURPLE)
    slide_label(s, "UNDER THE HOOD", PURPLE)
    section_title(s, "What a Vector Actually Is", size=28, y=Inches(0.65))
    divider(s, Inches(1.46))

    # Left — vector visual
    box(s, Inches(0.42), Inches(1.65), Inches(5.5), Inches(5.0), PANEL, BORDER, rounding=True)
    t(s, "ONE STORED FUNCTION", Inches(0.65), Inches(1.82), Inches(4.5), Inches(0.35),
      size=9, bold=True, color=GREY)
    vector_text = ("[0.21, −0.54, 0.87, 0.13, −0.92,\n"
                   " 0.45, 0.71, −0.33, 0.58,  0.19,\n"
                   " ···\n"
                   " −0.23, 0.67, 0.42, −0.88, 0.15]")
    t(s, vector_text, Inches(0.62), Inches(2.22), Inches(5.0), Inches(1.45),
      size=12, color=PURPLE, italic=False)
    t(s, "← 1,536 floats total",
      Inches(0.62), Inches(3.62), Inches(4.8), Inches(0.38),
      size=10, color=GREY, italic=True)

    # What it encodes
    t(s, "Encodes ALL of:", Inches(0.65), Inches(4.05), Inches(4.8), Inches(0.38),
      size=11, bold=True, color=WHITE)
    for i, (item, c) in enumerate([
        ("Function name label",           BLUE),
        ("GPT plain-English summary",     AMBER),
        ("Raw code syntax & structure",   GREEN),
    ]):
        box(s, Inches(0.7), Inches(4.48 + i * 0.52), Inches(0.35), Inches(0.35), c, rounding=True)
        t(s, item, Inches(1.15), Inches(4.48 + i * 0.52), Inches(4.2), Inches(0.38),
          size=11, color=WHITE)

    # Right — cosine similarity
    box(s, Inches(6.18), Inches(1.65), Inches(6.6), Inches(5.0), PANEL, BORDER, rounding=True)
    t(s, "COSINE SIMILARITY", Inches(6.42), Inches(1.82), Inches(5.5), Inches(0.35),
      size=9, bold=True, color=GREY)
    t(s, "sim(A, B) =  A · B\n             |A| × |B|",
      Inches(6.42), Inches(2.2), Inches(5.5), Inches(1.0),
      size=18, bold=True, color=PURPLE)
    t(s, ("Measures the angle between two vectors.\n"
          "1.0 = identical direction (perfect match)\n"
          "0.0 = perpendicular (completely unrelated)\n\n"
          "Magnitude-independent — a 2-line function\n"
          "and a 50-line function can score equally."),
      Inches(6.42), Inches(3.3), Inches(5.8), Inches(2.1),
      size=11, color=GREY)

    t(s, "Why not Euclidean distance?",
      Inches(6.42), Inches(5.3), Inches(5.5), Inches(0.38),
      size=11, bold=True, color=AMBER)
    t(s, "Vector magnitude reflects text length, not meaning.\nCosine ignores magnitude — only direction matters.",
      Inches(6.42), Inches(5.65), Inches(5.8), Inches(0.72),
      size=10.5, color=GREY)


def slide_09_results(prs):
    s = new_slide(prs)
    morph(s)
    top_bar(s, GREEN)
    slide_label(s, "RESULTS", GREEN)
    section_title(s, "What It Actually Does", size=28, y=Inches(0.65))
    divider(s, Inches(1.46))

    # Big stat bubbles
    stats = [
        ("6/6",   "semantic queries\nfound  (0 by text search)", GREEN),
        ("0.92",  "MRR — correct result\nalmost always #1",       BLUE),
        ("80%",   "Precision@5 — 4 of 5\nreturns are relevant",   AMBER),
        ("82%",   "best cosine score\n(vs 46% without summary)",   PURPLE),
    ]
    sw = Inches(2.9)
    for i, (val, label, color) in enumerate(stats):
        x = Inches(0.42) + i * (sw + Inches(0.25))
        box(s, x, Inches(1.65), sw, Inches(2.3), PANEL, color, bw=1.2, rounding=True)
        t(s, val,   x, Inches(1.82), sw, Inches(1.1),
          size=52, bold=True, color=color, align=PP_ALIGN.CENTER)
        t(s, label, x, Inches(2.88), sw, Inches(0.85),
          size=10, color=GREY, align=PP_ALIGN.CENTER)

    # Query comparison
    box(s, Inches(0.42), Inches(4.15), Inches(12.5), Inches(3.0), PANEL, BORDER, rounding=True)
    t(s, "CONCEPTUAL QUERIES — NO KEYWORD OVERLAP WITH FUNCTION NAMES",
      Inches(0.65), Inches(4.28), Inches(10.0), Inches(0.35),
      size=9, bold=True, color=GREY)

    rows = [
        ("Query",                          "Traditional", "Smart Search",  "Score"),
        ('"fetch product from database"',  "0 results",   "getProductInfo()", "82%"),
        ('"check if user is authenticated"',"0 results",  "assert_logged_in()","74%"),
        ('"calculate total order price"',  "0 results",   "computeLineItemTotal()","68%"),
        ('"delete item from cart"',        "0 results",   "removeCartItem()",  "73%"),
    ]
    col_x = [Inches(0.6),  Inches(6.0), Inches(8.6), Inches(11.7)]
    col_w = [Inches(5.3), Inches(2.4), Inches(3.1), Inches(1.5)]
    for ri, row in enumerate(rows):
        yy = Inches(4.6 + ri * 0.46)
        cs = [GREY, RED, GREEN, AMBER] if ri == 0 else [WHITE, RED, GREEN, AMBER]
        bolds = [True, False, False, True] if ri > 0 else [True]*4
        for ci, (cell, cc, bold) in enumerate(zip(row, cs, bolds)):
            t(s, cell, col_x[ci], yy, col_w[ci], Inches(0.42),
              size=9.5 if ri == 0 else 10.5, bold=bold, color=cc)
        if ri == 0:
            box(s, Inches(0.55), Inches(5.02), Inches(12.3), Inches(0.02), DIM)


def slide_10_engineering(prs):
    s = new_slide(prs)
    morph(s)
    top_bar(s, BLUE)
    slide_label(s, "ENGINEERING", BLUE)
    section_title(s, "4 Key Technical Decisions", size=28, y=Inches(0.65))
    divider(s, Inches(1.46))

    decisions = [
        (BLUE,   "ThreadingHTTPServer",
                 ["Search + indexing run in parallel threads",
                  "One-line change — Python stdlib, no dependencies",
                  "Thread-safe: handlers share no mutable state"]),
        (GREEN,  "Two-Phase Batched Indexing",
                 ["Phase 1: scan all files (local, no network)",
                  "Phase 2: send in batches of 50 to backend",
                  "50× faster than sequential approach"]),
        (AMBER,  "/config Endpoint",
                 ["backend/config.py is single source of truth",
                  "MIN_AI_QUERY_LENGTH fetched by extension at startup",
                  "Frontend + backend always enforce the same value"]),
        (PURPLE, "Relative Path Chunk IDs",
                 ["ID format: src/auth.py::verify_token",
                  "Portable — works after folder moves",
                  "Consistent across developer machines"]),
    ]

    cw = Inches(5.85)
    ch = Inches(2.5)
    positions = [
        (Inches(0.42), Inches(1.65)),
        (Inches(6.55), Inches(1.65)),
        (Inches(0.42), Inches(4.38)),
        (Inches(6.55), Inches(4.38)),
    ]

    for (x, y), (color, title, items) in zip(positions, decisions):
        box(s, x, y, cw, ch, PANEL, BORDER, bw=0.7, rounding=True)
        box(s, x, y, cw, Inches(0.055), color)
        t(s, title, x + Inches(0.2), y + Inches(0.12), cw - Inches(0.3), Inches(0.45),
          size=14, bold=True, color=color)
        for i, item in enumerate(items):
            t(s, f"→  {item}", x + Inches(0.2), y + Inches(0.65 + i * 0.55),
              cw - Inches(0.3), Inches(0.5), size=10.5, color=GREY)


def slide_11_challenges(prs):
    s = new_slide(prs)
    morph(s)
    top_bar(s, RED)
    slide_label(s, "CHALLENGES", RED)
    section_title(s, "Problems Faced & How They Were Solved", size=26, y=Inches(0.65))
    divider(s, Inches(1.46))

    challenges = [
        ("Null IDs crashing Pinecone",
         "Filter at both TS + Python layers:\nid is string and id.length > 0"),
        ("30+ min first-run indexing",
         "Two-phase batching: 200 s → 4 s\nfor 100 functions (50× speedup)"),
        ("Class chunks polluting results",
         "Filter c.type !== 'class' before\nembedding — methods indexed, not wrappers"),
        ("Progress lost on crash",
         "Save index.json after every batch\n→ resume from last checkpoint"),
        ("Search blocked by indexing",
         "ThreadingHTTPServer: one thread per\nrequest, search + index run in parallel"),
        ("Config drift across layers",
         "/config endpoint: single number in\nconfig.py, all layers read from it"),
    ]

    cw = Inches(3.85)
    ch = Inches(2.1)
    positions = [
        (Inches(0.42), Inches(1.65)),
        (Inches(4.48), Inches(1.65)),
        (Inches(8.55), Inches(1.65)),
        (Inches(0.42), Inches(3.98)),
        (Inches(4.48), Inches(3.98)),
        (Inches(8.55), Inches(3.98)),
    ]

    for (x, y), (problem, solution) in zip(positions, challenges):
        box(s, x, y, cw, ch, PANEL, BORDER, bw=0.7, rounding=True)
        # Problem label
        box(s, x + Inches(0.15), y + Inches(0.14), Inches(0.35), Inches(0.35),
            RED, rounding=True)
        t(s, "!", x + Inches(0.15), y + Inches(0.14), Inches(0.35), Inches(0.35),
          size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        t(s, problem, x + Inches(0.6), y + Inches(0.14), cw - Inches(0.72), Inches(0.4),
          size=10.5, bold=True, color=WHITE)
        box(s, x + Inches(0.15), y + Inches(0.6), cw - Inches(0.3), Inches(0.02), DIM)
        t(s, solution, x + Inches(0.18), y + Inches(0.7), cw - Inches(0.3), Inches(1.2),
          size=10.5, color=GREY, wrap=True)


def slide_12_future(prs):
    s = new_slide(prs)
    morph(s)
    top_bar(s, AMBER)
    slide_label(s, "FUTURE WORK", AMBER)
    section_title(s, "What Comes Next", size=28, y=Inches(0.65))
    divider(s, Inches(1.46))

    items = [
        (BLUE,   "Remote Backend Hosting",
                 "Eliminate local Python setup. One config change (backendUrl) points extension at cloud."),
        (GREEN,  "Git Branch-Aware Indexing",
                 "Detect branch switches via VS Code Git API, compare changed files between branches."),
        (PURPLE, "Cross-Function Semantic Linking",
                 "Build call graph at index time. Boost scores of connected functions for flow queries."),
        (AMBER,  "Local Embedding Models",
                 "nomic-embed-code or CodeBERT via Ollama — offline/air-gapped use, only embedder.py changes."),
        (RED,    "Feedback-Based Ranking",
                 "Track clicks and skips locally. Adjust threshold and weights per developer over time."),
    ]

    for i, (color, title, body) in enumerate(items):
        y = Inches(1.65 + i * 1.02)
        # Number circle
        circle(s, Inches(0.92), y + Inches(0.38), Inches(0.32), color)
        t(s, str(i+1), Inches(0.6), y + Inches(0.06), Inches(0.64), Inches(0.64),
          size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Card
        box(s, Inches(1.42), y, Inches(11.4), Inches(0.85), PANEL, BORDER, bw=0.5, rounding=True)
        box(s, Inches(1.42), y, Inches(0.055), Inches(0.85), color)
        t(s, title, Inches(1.65), y + Inches(0.05), Inches(3.2), Inches(0.42),
          size=13, bold=True, color=color)
        t(s, body, Inches(5.05), y + Inches(0.08), Inches(7.6), Inches(0.68),
          size=10.5, color=GREY)


def slide_13_summary(prs):
    s = new_slide(prs)
    morph(s)

    # Background geometric accents
    box(s, Inches(0), Inches(5.5), W, Inches(2.0), PANEL)
    box(s, Inches(0), Inches(5.5), W, Inches(0.06), BLUE)

    # Coloured dots top-right
    for i, c in enumerate([BLUE, GREEN, AMBER, PURPLE]):
        circle(s, Inches(11.5 + i * 0.52), Inches(0.55), Inches(0.16), c)

    t(s, "What Was Built",
      Inches(0.55), Inches(0.55), Inches(9.0), Inches(0.6),
      size=13, bold=True, color=GREY)

    t(s, "Smart Search",
      Inches(0.55), Inches(1.1), Inches(11.0), Inches(1.3),
      size=62, bold=True, color=WHITE)

    achievements = [
        (BLUE,   "Semantic code search inside VS Code — finds functions by what they do"),
        (GREEN,  "Two-phase batching made first-run indexing 50× faster"),
        (AMBER,  "GPT summaries improved cosine scores from ~46% to ~82%"),
        (PURPLE, "MRR 0.92 — correct result is first in 5 of 6 conceptual queries"),
        (RED,    "Crash-safe, concurrent, configuration-consistent production architecture"),
    ]

    for i, (color, text) in enumerate(achievements):
        y = Inches(2.5 + i * 0.53)
        circle(s, Inches(0.75), y + Inches(0.2), Inches(0.14), color)
        t(s, text, Inches(1.05), y, Inches(11.5), Inches(0.48),
          size=12, color=WHITE)

    t(s, "Nawfal Jalloul  ·  Master's Research Project  ·  May 2026",
      Inches(0.55), Inches(5.82), Inches(9.0), Inches(0.5),
      size=12, color=GREY)
    t(s, "Thank you",
      Inches(9.5), Inches(5.7), Inches(3.5), Inches(0.85),
      size=28, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)


# ── Assemble ───────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_architecture(prs)
    slide_05_indexing(prs)
    slide_06_summaries(prs)
    slide_07_search(prs)
    slide_08_vectors(prs)
    slide_09_results(prs)
    slide_10_engineering(prs)
    slide_11_challenges(prs)
    slide_12_future(prs)
    slide_13_summary(prs)

    out = os.path.join(DIR, "smart_search_presentation.pptx")
    prs.save(out)
    print(f"Saved → {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
